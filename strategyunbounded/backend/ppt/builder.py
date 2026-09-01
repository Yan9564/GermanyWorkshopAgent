from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Colours
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x1F, 0x29, 0x37)
C_BODY = RGBColor(0x37, 0x41, 0x51)
C_LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)
C_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
C_ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
C_GREY_TEXT = RGBColor(0x6B, 0x72, 0x80)

PRIORITY_COLOURS = {
    1: RGBColor(0x22, 0xC5, 0x5E),
    2: RGBColor(0xF5, 0x9E, 0x0B),
    3: RGBColor(0x9C, 0xA3, 0xAF),
}
PRIORITY_LABELS = {1: "Priority 1", 2: "Priority 2", 3: "Priority 3"}


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE
    return slide


def _add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                 color=None, align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_DARK
    return txb


def _add_badge(slide, left, top, width, height, label, bg_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = C_WHITE


def _add_divider(slide, top):
    shape = slide.shapes.add_shape(
        1,
        Inches(0.5), top,
        Inches(12.33), Pt(1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BORDER
    shape.line.fill.background()


def _add_label_value(slide, left, top, width, label, value, label_size=9, value_size=11):
    _add_textbox(slide, left, top, width, Inches(0.3), label.upper(),
                 font_size=label_size, bold=True, color=C_GREY_TEXT)
    _add_textbox(slide, left, top + Inches(0.3), width, Inches(0.8), value,
                 font_size=value_size, color=C_BODY, word_wrap=True)


def _slide_title(prs: Presentation) -> None:
    slide = _blank_slide(prs)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1D, 0x4E, 0xD8)

    _add_textbox(
        slide,
        Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
        "AI Use Case Workshop Results",
        font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
        "Service Workshop — Part 1",
        font_size=20, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER,
    )


def _slide_problems(prs: Presentation, problems: list[dict]) -> None:
    slide = _blank_slide(prs)

    _add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                 "Your 5 Problems", font_size=28, bold=True)
    _add_divider(slide, Inches(1.3))

    for i, prob in enumerate(problems):
        top = Inches(1.5) + i * Inches(1.1)
        _add_textbox(slide, Inches(0.5), top, Inches(0.6), Inches(0.8),
                     f"{i + 1}.", font_size=16, bold=True, color=C_ACCENT)
        text = prob["text"] if len(prob["text"]) <= 200 else prob["text"][:197] + "…"
        _add_textbox(slide, Inches(1.1), top, Inches(11.7), Inches(0.9),
                     text, font_size=14, color=C_BODY)


def _slide_top_priorities(prs: Presentation, problems: list[dict]) -> None:
    slide = _blank_slide(prs)

    _add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                 "Your Top Priorities", font_size=28, bold=True)
    _add_divider(slide, Inches(1.3))

    all_uc: dict[int, dict | None] = {1: None, 2: None, 3: None}
    problem_texts: dict[str, str] = {p["id"]: p["text"] for p in problems}

    for prob in problems:
        for uc in prob.get("useCases", []):
            up = uc.get("userPriority")
            if up in (1, 2, 3):
                all_uc[up] = {**uc, "_problem_text": prob["text"], "_problem_index": prob["index"]}

    for slot in (1, 2, 3):
        top = Inches(1.5) + (slot - 1) * Inches(1.8)
        uc = all_uc[slot]
        color = PRIORITY_COLOURS[slot]

        _add_badge(slide, Inches(0.5), top + Inches(0.1), Inches(1.4), Inches(0.4),
                   PRIORITY_LABELS[slot], color)

        if uc:
            _add_textbox(slide, Inches(2.1), top, Inches(10), Inches(0.5),
                         uc["title"], font_size=16, bold=True)
            prob_label = f"↳ Problem {uc['_problem_index'] + 1}: {uc['_problem_text'][:80]}…" \
                if len(uc["_problem_text"]) > 80 \
                else f"↳ Problem {uc['_problem_index'] + 1}: {uc['_problem_text']}"
            _add_textbox(slide, Inches(2.1), top + Inches(0.55), Inches(10), Inches(0.4),
                         prob_label, font_size=11, color=C_GREY_TEXT, italic=True)
        else:
            _add_textbox(slide, Inches(2.1), top + Inches(0.1), Inches(10), Inches(0.5),
                         "Not assigned", font_size=14, color=C_GREY_TEXT, italic=True)


def _slide_problem_header(prs: Presentation, problem: dict) -> None:
    slide = _blank_slide(prs)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
                 f"Problem {problem['index'] + 1}", font_size=14, bold=True, color=C_ACCENT)

    text = problem["text"]
    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5),
                 text, font_size=24, bold=True, color=C_DARK)


def _slide_use_case(prs: Presentation, uc: dict, problem: dict) -> None:
    slide = _blank_slide(prs)

    top_offset = Inches(0.35)

    user_priority = uc.get("userPriority")
    if user_priority in PRIORITY_COLOURS:
        _add_badge(slide, Inches(0.5), top_offset, Inches(1.4), Inches(0.38),
                   PRIORITY_LABELS[user_priority], PRIORITY_COLOURS[user_priority])

    feedback = uc.get("feedback")
    if feedback == "up":
        _add_textbox(slide, Inches(12.3), top_offset, Inches(0.6), Inches(0.4),
                     "👍", font_size=18, align=PP_ALIGN.RIGHT)
    elif feedback == "down":
        _add_textbox(slide, Inches(12.3), top_offset, Inches(0.6), Inches(0.4),
                     "👎", font_size=18, align=PP_ALIGN.RIGHT)

    _add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.7),
                 uc["title"], font_size=26, bold=True)
    _add_textbox(slide, Inches(0.5), Inches(1.55), Inches(12.33), Inches(0.45),
                 uc["summary"], font_size=13, italic=True, color=C_BODY)

    _add_divider(slide, Inches(2.1))

    _add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.33), Inches(0.7),
                 uc["description"], font_size=12, color=C_BODY)

    how_it_works = uc.get("howItWorks", [])
    hiw_top = Inches(3.0)
    _add_textbox(slide, Inches(0.5), hiw_top, Inches(4.5), Inches(0.3),
                 "HOW IT WORKS", font_size=9, bold=True, color=C_GREY_TEXT)
    for j, step in enumerate(how_it_works[:5]):
        _add_textbox(slide, Inches(0.5), hiw_top + Inches(0.35) + j * Inches(0.42),
                     Inches(4.5), Inches(0.4), f"• {step}", font_size=11, color=C_BODY)

    right_col = Inches(5.5)
    col_w = Inches(3.7)

    _add_label_value(slide, right_col, Inches(3.0), col_w, "Data Required", uc["dataRequired"])
    _add_label_value(slide, right_col, Inches(4.2), col_w, "Time to Implement", uc["timeToImplement"])

    right_col2 = Inches(9.4)
    _add_label_value(slide, right_col2, Inches(3.0), col_w, "Complexity", uc["complexity"])
    _add_label_value(slide, right_col2, Inches(4.2), col_w, "Est. Cost / ROI", uc["estimatedCostRoi"])

    _add_textbox(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
                 f"Problem {problem['index'] + 1}", font_size=9, color=C_GREY_TEXT)


def build_pptx(session: dict) -> BytesIO:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs)
    _slide_problems(prs, session["problems"])
    _slide_top_priorities(prs, session["problems"])

    for problem in session["problems"]:
        _slide_problem_header(prs, problem)
        for uc in problem.get("useCases", []):
            _slide_use_case(prs, uc, problem)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
